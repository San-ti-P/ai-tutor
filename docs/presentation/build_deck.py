#!/usr/bin/env python3
"""
'Avances de Implementacion' deck — AI Tutor TP2.
Evidence-first: real LangGraph agent graphs + a real Langfuse exam-generation
trace + an honest epic/branch status. No vanity code/trace metrics.

Assets (generated, real):
  assets/graph_*.png   -> rendered from the dev code (render_graphs.py)
  assets/exam_trace.png-> reconstructed from traces.csv (build_exam_trace.py)
"""
import json
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BASE = "/home/santiago/workspaces/ai-tutor/docs"
ASSETS = f"{BASE}/presentation/assets"
TRACE_META = json.load(open(f"{BASE}/observability/exam_trace_meta.json"))

# ---- Palette ----
INDIGO = RGBColor(0x4F, 0x46, 0xE5)
DARK   = RGBColor(0x1E, 0x1B, 0x4B)
ACCENT = RGBColor(0x8B, 0x5C, 0xF6)
VIOLET = RGBColor(0x6D, 0x28, 0xD9)
LIGHT  = RGBColor(0xF5, 0xF3, 0xFF)
CARD   = RGBColor(0xEE, 0xEC, 0xFB)
TEXT   = RGBColor(0x1F, 0x29, 0x37)
MUTED  = RGBColor(0x6B, 0x72, 0x80)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREEN  = RGBColor(0x10, 0xB9, 0x81)
AMBER  = RGBColor(0xF5, 0x9E, 0x0B)
FONT = "Calibri"; FONT_L = "Calibri Light"

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)

def rect(s, x, y, w, h, fill, line=None, shape=MSO_SHAPE.RECTANGLE):
    sp = s.shapes.add_shape(shape, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = line; sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp

def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=4, line_spacing=1.0):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Pt(2); tf.margin_right = Pt(2); tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    paras = runs if runs and isinstance(runs[0], list) else [runs]
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space_after); p.line_spacing = line_spacing
        for (txt, size, bold, color, *fn) in para:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
            r.font.name = fn[0] if fn else FONT
    return tb

def header(s, kicker, title, idx):
    rect(s, 0, 0, SW, Inches(1.2), DARK)
    rect(s, 0, Inches(1.2), SW, Pt(4), ACCENT)
    text(s, Inches(0.55), Inches(0.16), Inches(11), Inches(0.32),
         [[(kicker.upper(), 12, True, ACCENT)]])
    text(s, Inches(0.55), Inches(0.46), Inches(11.6), Inches(0.66),
         [[(title, 25, True, WHITE)]])
    text(s, Inches(12.4), Inches(0.42), Inches(0.7), Inches(0.5),
         [[(f"{idx:02d}", 14, True, ACCENT)]], align=PP_ALIGN.RIGHT)
    text(s, Inches(0.55), Inches(7.09), Inches(9), Inches(0.3),
         [[("Tutor Académico Personal · IA 2026 · UTN Santa Fe — CIDISI", 9, False, MUTED)]])

def picture_in_col(s, path, cx, top, height):
    """Add picture at given HEIGHT, horizontally centered on cx."""
    pic = s.shapes.add_picture(path, Emu(0), top, height=height)
    pic.left = Emu(int(cx - pic.width / 2))
    return pic


# ===========================================================================
# 1 — TITLE
# ===========================================================================
s = slide()
rect(s, 0, 0, SW, SH, DARK)
rect(s, 0, Inches(4.5), SW, Pt(3), ACCENT)
rect(s, Inches(0.9), Inches(1.0), Inches(1.9), Inches(0.5), ACCENT)
text(s, Inches(0.9), Inches(1.02), Inches(3), Inches(0.45),
     [[("ENTREGA 2 · AVANCES", 13, True, WHITE)]], anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
text(s, Inches(0.9), Inches(2.05), Inches(11.5), Inches(1.4),
     [[("Tutor Académico Personal", 50, True, WHITE)]])
text(s, Inches(0.9), Inches(3.35), Inches(11.5), Inches(0.8),
     [[("Avances de implementación — agentes en ejecución", 23, False, ACCENT, FONT_L)]])
text(s, Inches(0.9), Inches(4.85), Inches(11.5), Inches(1.0),
     [[("Grafos reales de cada agente · traza real de generación de examen · estado por épica",
        16, False, RGBColor(0xC4, 0xB5, 0xFD))]])
text(s, Inches(0.9), Inches(6.45), Inches(11.5), Inches(0.5),
     [[("Trabajo Práctico N°2  ·  Inteligencia Artificial 2026  ·  UTN Santa Fe — CIDISI",
        14, False, RGBColor(0x9C, 0xA3, 0xAF))]])

# ===========================================================================
# 2 — ÉPICAS: cómo avanzamos (honest status)
# ===========================================================================
s = slide()
header(s, "Cómo avanzamos", "Una épica por rama, integradas a dev por PR", 2)
text(s, Inches(0.55), Inches(1.4), Inches(12.2), Inches(0.5),
     [[("Cada épica se desarrolló en su rama ", 14, False, TEXT),
       ("feat/epic-XX", 14, True, INDIGO),
       (" y se integró a ", 14, False, TEXT),
       ("dev", 14, True, INDIGO),
       (" mediante Pull Request.", 14, False, TEXT)]])

rows = [
    ("01", "Orchestrator", "dev", True),
    ("02", "Ingestor", "dev", True),
    ("03", "ExamGenerator", "dev", True),
    ("04", "ExerciseGenerator", "dev", True),
    ("05", "Evaluator", "dev", True),
    ("06", "Support Agent", "dev", True),
    ("07", "UI — Next.js", "dev", True),
    ("10", "Refactor de arquitectura", "dev", True),
    ("08", "Observabilidad (Langfuse)", "feat/epic-08 · en progreso", False),
    ("09", "Profile bootstrap", "identificada · gap a cerrar", None),
]
table = s.shapes.add_table(len(rows) + 1, 3, Inches(0.55), Inches(2.0),
                           Inches(12.2), Inches(4.7)).table
table.columns[0].width = Inches(1.4)
table.columns[1].width = Inches(5.6)
table.columns[2].width = Inches(5.2)
for c, h in enumerate(["Épica", "Componente", "Estado"]):
    cell = table.cell(0, c); cell.fill.solid(); cell.fill.fore_color.rgb = INDIGO
    cell.text = h; p = cell.text_frame.paragraphs[0]
    p.font.bold = True; p.font.size = Pt(13); p.font.color.rgb = WHITE; p.font.name = FONT
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
for r, (num, comp, state, done) in enumerate(rows, start=1):
    badge = "✅  en dev" if done else ("⏳  " + state if done is False else "📋  " + state)
    state_txt = ("en dev" if done else state)
    vals = [num, comp, ("✅  en dev" if done else ("⏳  " if done is False else "📋  ") + state)]
    for c, v in enumerate(vals):
        cell = table.cell(r, c); cell.fill.solid()
        cell.fill.fore_color.rgb = WHITE if r % 2 else LIGHT
        cell.text = v; p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(12.5); p.font.name = FONT
        p.font.bold = (c == 0)
        if c == 2 and done:
            p.font.color.rgb = GREEN; p.font.bold = True
        elif c == 2 and done is False:
            p.font.color.rgb = AMBER; p.font.bold = True
        elif c == 2:
            p.font.color.rgb = MUTED
        else:
            p.font.color.rgb = INDIGO if c == 0 else TEXT
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

# ===========================================================================
# 3 — ARQUITECTURA multi-agente (brief)
# ===========================================================================
s = slide()
header(s, "Arquitectura", "Orquestador + 5 agentes especializados", 3)
# central hub
rect(s, Inches(5.3), Inches(2.0), Inches(2.7), Inches(1.1), INDIGO, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, Inches(5.3), Inches(2.0), Inches(2.7), Inches(1.1),
     [[("Orchestrator", 16, True, WHITE)], [("Plan-and-Execute", 10.5, False, RGBColor(0xC4,0xB5,0xFD))]],
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
spokes = [
    ("Ingestor", "ingiere material → RAG"),
    ("ExamGenerator", "exámenes MCQ + abiertas"),
    ("ExerciseGenerator", "ejercicios prácticos"),
    ("Evaluator", "corrige + feedback"),
    ("Support", "perfil y progreso"),
]
positions = [
    (Inches(0.7), Inches(4.0)), (Inches(3.2), Inches(4.0)),
    (Inches(5.7), Inches(4.0)), (Inches(8.2), Inches(4.0)),
    (Inches(10.7), Inches(4.0)),
]
for (name, desc), (x, y) in zip(spokes, positions):
    rect(s, x, y, Inches(2.3), Inches(1.5), WHITE, line=CARD, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, x, y, Inches(2.3), Inches(0.12), ACCENT)
    text(s, x + Inches(0.1), y + Inches(0.3), Inches(2.1), Inches(0.5),
         [[(name, 13.5, True, INDIGO)]], align=PP_ALIGN.CENTER)
    text(s, x + Inches(0.1), y + Inches(0.82), Inches(2.1), Inches(0.6),
         [[(desc, 10.5, False, MUTED)]], align=PP_ALIGN.CENTER, line_spacing=0.95)
    # connector
    conn = rect(s, x + Inches(1.05), Inches(3.1), Pt(2), Inches(0.9), CARD)

text(s, Inches(0.55), Inches(6.05), Inches(12.2), Inches(0.9),
     [[("Cada agente es un ", 13, False, TEXT),
       ("StateGraph de LangGraph", 13, True, VIOLET),
       (" con su propio flujo de control. El Orchestrator clasifica la intención del "
        "usuario y enruta al agente correcto (o planifica tareas compuestas).", 13, False, TEXT)]],
     line_spacing=1.1)

# ===========================================================================
# 4 — GRAFO REAL: Orchestrator
# ===========================================================================
s = slide()
header(s, "Grafo real · 1/4", "Orchestrator — enrutamiento y plan-execute", 4)
picture_in_col(s, f"{ASSETS}/graph_orchestrator.png", cx=Inches(3.4), top=Inches(1.55), height=Inches(5.2))
rect(s, Inches(6.8), Inches(1.7), Inches(5.95), Inches(4.9), LIGHT, line=CARD)
text(s, Inches(7.05), Inches(1.9), Inches(5.5), Inches(0.5),
     [[("Qué hace el grafo", 15, True, VIOLET)]])
nodes = [
    ("classify_intent", "detecta la intención (ingest, examen, ejercicio, evaluar, perfil, compuesto)"),
    ("route_to_agent", "arista condicional → deriva al agente adecuado"),
    ("plan_composite", "descompone pedidos multi-paso en una secuencia de tools"),
    ("execute_step", "ejecuta cada paso; bucle con control de iteraciones"),
    ("synthesize_response", "redacta la respuesta final al estudiante"),
]
tb = s.shapes.add_textbox(Inches(7.05), Inches(2.45), Inches(5.5), Inches(4.2))
tf = tb.text_frame; tf.word_wrap = True
for i, (h, b) in enumerate(nodes):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(11)
    r = p.add_run(); r.text = h + "\n"; r.font.bold = True; r.font.size = Pt(13.5); r.font.color.rgb = TEXT
    r = p.add_run(); r.text = b; r.font.size = Pt(12); r.font.color.rgb = MUTED
text(s, Inches(0.55), Inches(6.78), Inches(6), Inches(0.3),
     [[("Renderizado del código real en dev (LangGraph .draw_mermaid)", 9.5, False, MUTED)]])

# ===========================================================================
# 5 — GRAFOS REALES: generadores (Ingestor, Exam, Exercise)
# ===========================================================================
s = slide()
header(s, "Grafo real · 2/4", "Ingestor · ExamGenerator · ExerciseGenerator", 5)
cols = [
    ("Ingestor", "graph_ingestor.png", "lineal: parse → classify → chunk+embed"),
    ("ExamGenerator", "graph_exam_generator.png", "retrieve → generate → validate → format"),
    ("ExerciseGenerator", "graph_exercise_generator.png", "genera ejercicios con contexto RAG"),
]
centers = [Inches(2.55), Inches(6.67), Inches(10.78)]
for (title, img, desc), cx in zip(cols, centers):
    text(s, cx - Inches(2.0), Inches(1.45), Inches(4.0), Inches(0.4),
         [[(title, 16, True, INDIGO)]], align=PP_ALIGN.CENTER)
    picture_in_col(s, f"{ASSETS}/{img}", cx=cx, top=Inches(2.0), height=Inches(4.0))
    text(s, cx - Inches(2.0), Inches(6.2), Inches(4.0), Inches(0.7),
         [[(desc, 11.5, False, MUTED)]], align=PP_ALIGN.CENTER, line_spacing=0.95)

# ===========================================================================
# 6 — GRAFOS REALES: Evaluator + Support (los más complejos)
# ===========================================================================
s = slide()
header(s, "Grafo real · 3/4", "Evaluator y Support — ruteo condicional", 6)
# Evaluator (tall) on left
text(s, Inches(0.4), Inches(1.4), Inches(4.6), Inches(0.4),
     [[("Evaluator", 16, True, INDIGO)]], align=PP_ALIGN.CENTER)
picture_in_col(s, f"{ASSETS}/graph_evaluator.png", cx=Inches(2.7), top=Inches(1.85), height=Inches(5.0))
# Support
text(s, Inches(5.3), Inches(1.4), Inches(3.0), Inches(0.4),
     [[("Support", 16, True, INDIGO)]], align=PP_ALIGN.CENTER)
picture_in_col(s, f"{ASSETS}/graph_support.png", cx=Inches(6.6), top=Inches(2.1), height=Inches(3.6))
# notes panel
rect(s, Inches(8.55), Inches(1.7), Inches(4.25), Inches(4.95), LIGHT, line=CARD)
text(s, Inches(8.8), Inches(1.9), Inches(3.8), Inches(0.5),
     [[("Ingeniería de control", 14, True, VIOLET)]])
notes = [
    ("Evaluator (8 nodos)", "ramifica según evaluabilidad y dispara un llm_judge sólo cuando hace falta validar el feedback"),
    ("Anti-alucinación", "el feedback se valida contra el material antes de entregarse"),
    ("Support (4 nodos)", "salta la historia de sesiones si el alumno es nuevo (arista condicional)"),
]
tb = s.shapes.add_textbox(Inches(8.8), Inches(2.45), Inches(3.85), Inches(4.0))
tf = tb.text_frame; tf.word_wrap = True
for i, (h, b) in enumerate(notes):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(12)
    r = p.add_run(); r.text = h + "\n"; r.font.bold = True; r.font.size = Pt(13); r.font.color.rgb = TEXT
    r = p.add_run(); r.text = b; r.font.size = Pt(11.5); r.font.color.rgb = MUTED

# ===========================================================================
# 7 — TRAZA REAL: generación de examen
# ===========================================================================
s = slide()
header(s, "Grafo real · 4/4 · ejecución", "Traza real: generación de un examen", 7)
picture_in_col(s, f"{ASSETS}/exam_trace.png", cx=Inches(4.6), top=Inches(1.4), height=Inches(5.5))
rect(s, Inches(9.3), Inches(1.55), Inches(3.5), Inches(5.1), LIGHT, line=CARD)
text(s, Inches(9.55), Inches(1.75), Inches(3.1), Inches(0.5),
     [[("Captura real (Langfuse)", 13.5, True, VIOLET)]])
pts = [
    "Un pedido real disparó el flujo completo de extremo a extremo.",
    "classify_intent → plan_composite → execute_step → generate_exam.",
    "generate_exam abre su propio sub-grafo LangGraph.",
    "Adentro: retrieve_chunks → rag_retrieve (RAG real sobre ChromaDB).",
    f"{TRACE_META['n_observations_in_trace']} observaciones trazadas en esta sola corrida.",
]
tb = s.shapes.add_textbox(Inches(9.55), Inches(2.35), Inches(3.05), Inches(4.0))
tf = tb.text_frame; tf.word_wrap = True
for i, p_ in enumerate(pts):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(10)
    r = p.add_run(); r.text = "• "; r.font.bold = True; r.font.size = Pt(12); r.font.color.rgb = ACCENT
    r = p.add_run(); r.text = p_; r.font.size = Pt(11.5); r.font.color.rgb = TEXT
text(s, Inches(0.55), Inches(6.95), Inches(8), Inches(0.3),
     [[("Reconstruida del export de trazas · modelo local Gemma (Ollama)", 9.5, False, MUTED)]])

# ===========================================================================
# 8 — ESTADO HONESTO + PRÓXIMOS PASOS
# ===========================================================================
s = slide()
header(s, "Estado y próximos pasos", "Qué está listo y qué sigue", 8)
# done
rect(s, Inches(0.55), Inches(1.6), Inches(6.0), Inches(4.4), WHITE, line=CARD)
rect(s, Inches(0.55), Inches(1.6), Inches(6.0), Inches(0.6), GREEN)
text(s, Inches(0.55), Inches(1.6), Inches(6.0), Inches(0.6),
     [[("✓  EN dev — FUNCIONA", 15, True, WHITE)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
done = [
    "6 agentes en LangGraph, cada uno con su grafo",
    "Orquestación hub-and-spoke + RAG sobre ChromaDB",
    "Memoria SQLite + API FastAPI + UI Next.js",
    "Refactor de arquitectura (épica 10)",
]
tb = s.shapes.add_textbox(Inches(0.9), Inches(2.45), Inches(5.4), Inches(3.4))
tf = tb.text_frame; tf.word_wrap = True
for i, it in enumerate(done):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(13)
    r = p.add_run(); r.text = "✓  "; r.font.bold = True; r.font.size = Pt(14); r.font.color.rgb = GREEN
    r = p.add_run(); r.text = it; r.font.size = Pt(13.5); r.font.color.rgb = TEXT
# in progress
rect(s, Inches(6.75), Inches(1.6), Inches(6.0), Inches(4.4), WHITE, line=CARD)
rect(s, Inches(6.75), Inches(1.6), Inches(6.0), Inches(0.6), AMBER)
text(s, Inches(6.75), Inches(1.6), Inches(6.0), Inches(0.6),
     [[("▣  EN PROGRESO / SIGUE", 15, True, WHITE)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
todo = [
    "Épica 08 · Observabilidad: integrar Langfuse a dev (hoy en su rama)",
    "Épica 09 · Profile bootstrap: cerrar el gap de carga de perfil",
    "OCR de matemática (fórmulas → LaTeX)",
    "Coloquio y demo en vivo · entrega final 29/06",
]
tb = s.shapes.add_textbox(Inches(7.1), Inches(2.45), Inches(5.4), Inches(3.4))
tf = tb.text_frame; tf.word_wrap = True
for i, it in enumerate(todo):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(13)
    r = p.add_run(); r.text = "○  "; r.font.bold = True; r.font.size = Pt(14); r.font.color.rgb = AMBER
    r = p.add_run(); r.text = it; r.font.size = Pt(13.5); r.font.color.rgb = TEXT
text(s, Inches(0.55), Inches(6.25), Inches(12.2), Inches(0.7),
     [[("Las trazas que mostramos se generaron ejecutando la rama de observabilidad; "
        "el próximo paso es integrarla a dev.", 12, False, MUTED)]])

# ===========================================================================
# 9 — CIERRE
# ===========================================================================
s = slide()
rect(s, 0, 0, SW, SH, DARK)
rect(s, Inches(0.9), Inches(2.6), Inches(1.9), Pt(5), ACCENT)
text(s, Inches(0.9), Inches(2.85), Inches(11.5), Inches(1.3),
     [[("Seis agentes que ejecutan de verdad", 38, True, WHITE)]])
text(s, Inches(0.9), Inches(4.15), Inches(11.5), Inches(0.8),
     [[("Grafos reales del código · una traza real de punta a punta · estado honesto por épica",
        17, False, RGBColor(0xC4, 0xB5, 0xFD))]])
text(s, Inches(0.9), Inches(5.5), Inches(11.5), Inches(0.6),
     [[("¿Pasamos a la demo en vivo?", 20, True, ACCENT)]])
text(s, Inches(0.9), Inches(6.55), Inches(11.5), Inches(0.5),
     [[("Tutor Académico Personal · IA 2026 · UTN Santa Fe — CIDISI", 13, False, RGBColor(0x9C,0xA3,0xAF))]])

OUT = f"{BASE}/presentation/Avances_Tutor_Academico_TP2.pptx"
prs.save(OUT)
print(f"✅ Guardado: {OUT}")
print(f"   Slides: {len(prs.slides._sldIdLst)}")
